#!/usr/bin/env python3
"""
PowerPoint Content Extractor

Extracts PowerPoint presentation content to text/markdown format for use in LLM-based content synthesis.

Primary use case: Course material presentations -> readable text for synthesis

Usage:
    python pptx_extract.py presentation.pptx                    # Extract to markdown
    python pptx_extract.py presentation.pptx -f txt             # Extract to plain text
    python pptx_extract.py presentation.pptx --notes            # Include speaker notes
    python pptx_extract.py --dir /path/to/slides -o ./out       # Batch process directory
"""

import argparse
import sys
from pathlib import Path
from typing import Optional

try:
    from markitdown import MarkItDown
except ImportError:
    print("ERROR: markitdown is required. Install with: pip install markitdown", file=sys.stderr)
    sys.exit(1)


def extract_with_markitdown(pptx_path: Path) -> str:
    """Extract presentation content using markitdown."""
    md = MarkItDown()
    result = md.convert(str(pptx_path))
    return result.text_content


def extract_with_python_pptx(pptx_path: Path, include_notes: bool = False) -> str:
    """Extract presentation content using python-pptx (fallback method with notes support)."""
    try:
        from pptx import Presentation
    except ImportError:
        print("ERROR: python-pptx is required for --notes. Install with: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(pptx_path))
    slides_content = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []

        # Extract title if present
        if slide.shapes.title:
            title = slide.shapes.title.text
            slide_text.append(f"## Slide {slide_num}: {title}")
        else:
            slide_text.append(f"## Slide {slide_num}")

        slide_text.append("")

        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Skip if it's the title we already added
                if shape == slide.shapes.title:
                    continue

                text = shape.text.strip()
                # Format as bullet points if multiple lines
                lines = text.split("\n")
                for line in lines:
                    if line.strip():
                        slide_text.append(f"- {line.strip()}")

        # Extract speaker notes if requested
        if include_notes and slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                slide_text.append("")
                slide_text.append(f"> **Speaker Notes:** {notes_frame.text.strip()}")

        slides_content.append("\n".join(slide_text))

    return "\n\n".join(slides_content)


def process_file(pptx_path: Path, output_dir: Path, fmt: str, include_notes: bool) -> Optional[Path]:
    """Process a single pptx file."""
    try:
        # Use python-pptx if notes are requested, otherwise use markitdown
        if include_notes:
            text = extract_with_python_pptx(pptx_path, include_notes=True)
        else:
            try:
                text = extract_with_markitdown(pptx_path)
            except Exception:
                # Fallback to python-pptx if markitdown fails
                text = extract_with_python_pptx(pptx_path, include_notes=False)

        if not text.strip():
            return None

        # Format output
        ext = ".md" if fmt == "md" else ".txt"
        output_path = output_dir / f"{pptx_path.stem}{ext}"

        if fmt == "md":
            content = f"# {pptx_path.stem}\n\n{text}"
        else:
            # Strip markdown formatting for plain text
            content = text

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(content)

        return output_path

    except Exception as e:
        print(f"  WARNING: Failed to extract '{pptx_path.name}': {e}", file=sys.stderr)
        return None


def main():
    parser = argparse.ArgumentParser(
        description="Extract PowerPoint presentation content to text/markdown format",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Extract to markdown (default)
  python pptx_extract.py presentation.pptx

  # Extract to plain text
  python pptx_extract.py presentation.pptx -f txt

  # Include speaker notes
  python pptx_extract.py presentation.pptx --notes

  # Batch process directory
  python pptx_extract.py --dir /path/to/slides -o ./extracted
        """
    )

    parser.add_argument("presentation", nargs="?", help="Path to PowerPoint file (.pptx)")
    parser.add_argument("-o", "--output",
                        help="Output file or directory")
    parser.add_argument("-f", "--format",
                        choices=["md", "txt"],
                        default="md",
                        help="Output format (default: md)")
    parser.add_argument("-d", "--dir",
                        help="Process directory of presentations")
    parser.add_argument("-r", "--recursive", action="store_true",
                        help="Recursive directory scan")
    parser.add_argument("--notes", action="store_true",
                        help="Include speaker notes")
    parser.add_argument("-q", "--quiet", action="store_true",
                        help="Suppress progress output")

    args = parser.parse_args()

    def log(msg):
        if not args.quiet:
            print(msg)

    # Determine input files
    if args.dir:
        input_dir = Path(args.dir)
        if not input_dir.exists():
            print(f"ERROR: Directory not found: {input_dir}", file=sys.stderr)
            sys.exit(1)

        pattern = "**/*.pptx" if args.recursive else "*.pptx"
        pptx_files = list(input_dir.glob(pattern))

        if not pptx_files:
            print("No .pptx files found", file=sys.stderr)
            sys.exit(1)

        log(f"Found {len(pptx_files)} presentations")

        # Determine output directory
        output_dir = Path(args.output) if args.output else input_dir
        output_dir.mkdir(parents=True, exist_ok=True)

        # Process all files
        exported = 0
        failed = 0

        for pptx_path in pptx_files:
            result = process_file(pptx_path, output_dir, args.format, args.notes)

            if result:
                log(f"  {pptx_path.name} -> {result.name}")
                exported += 1
            else:
                log(f"  {pptx_path.name} - FAILED")
                failed += 1

        log(f"\nExtracted {exported} presentations to: {output_dir}")
        if failed:
            log(f"Failed: {failed} presentations")

    elif args.presentation:
        pptx_path = Path(args.presentation)

        if not pptx_path.exists():
            print(f"ERROR: File not found: {pptx_path}", file=sys.stderr)
            sys.exit(1)

        # Determine output path
        if args.output:
            output_path = Path(args.output)
            if output_path.suffix in [".md", ".txt"]:
                output_dir = output_path.parent
            else:
                output_dir = output_path
        else:
            output_dir = pptx_path.parent

        output_dir.mkdir(parents=True, exist_ok=True)

        log(f"Extracting: {pptx_path.name}")
        result = process_file(pptx_path, output_dir, args.format, args.notes)

        if result:
            log(f"Output: {result}")
        else:
            print("ERROR: Extraction failed", file=sys.stderr)
            sys.exit(1)

    else:
        parser.print_help()
        sys.exit(1)


if __name__ == "__main__":
    main()
